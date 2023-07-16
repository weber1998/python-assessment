from pptx import Presentation
from PIL import Image
class ToPptx:
    def __init__(self, json_data):
        self.json_data = json_data
        self.presentation = Presentation()

    def presentation_build(self):
        for slide_data in self.json_data:
            slide_type = slide_data['type']
            slide_title = slide_data['title']
            slide_content = slide_data['content']

            slide_layout = self._get_slide_layout(slide_type)
            slide = self.presentation.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = slide_title

            if slide_type == 'title':
                subtitle = slide.placeholders[1]
                subtitle.text = slide_content

            elif slide_type == 'text':
                content = slide.placeholders[1]
                content.text = slide_content

            elif slide_type == 'list':
                content_placeholder = slide.placeholders[1]
                for item in slide_content:
                    level = item['level']
                    text = item['text']
                    content_placeholder.text += f"\n{' ' * (level * 4)}- {text}"

            elif slide_type == 'picture':
                image_path = 'picture.png'
                with open(image_path, 'rb') as img_file:
                    img_file.read()
                try:
                    image = Image.open(image_path)
                except FileNotFoundError:
                    print(f"Image file '{image_path}' not found.")
                except Exception as e:
                    print(f"Error occurred while opening the image: {str(e)}")

                slide_width = self.presentation.slide_width
                slide_height = self.presentation.slide_height

                image_width = image.width * 12700
                image_height = image.height * 12700

                # középre igaztás (nem működik megfelelően)
                img_half_width = image_width / 2
                img_half_height = image_height / 2
                left = (slide_width / 2) - (img_half_width / 2)
                top = (slide_height / 2) - (img_half_height / 2)

                slide.shapes.add_picture(slide_content, left, top)

            elif slide_type == 'plot':
                content_placeholder = slide.placeholders[0]
                content_placeholder.text = slide_title
                configuration = slide_data.get('configuration', {})
                x_label = configuration.get('x-label', '')
                y_label = configuration.get('y-label', '')

    def _get_slide_layout(self, slide_type):
        slide_layouts = self.presentation.slide_layouts
        if slide_type == 'title':
            return slide_layouts[0]
        elif slide_type == 'text':
            return slide_layouts[1]
        elif slide_type == 'list':
            return slide_layouts[1]
        elif slide_type == 'picture':
            return slide_layouts[5]
        elif slide_type == 'plot':
            return slide_layouts[5]

    def save(self, filename):
        self.presentation.save(filename)